#!/usr/bin/env python3
"""
extract_office_doc.py
Gitea Office Document Text Extractor

Downloads a file from Gitea (or reads local file) and extracts text content
from PDF / Word / PowerPoint / Excel files.

Usage:
  python3 extract_office_doc.py --url <gitea_raw_url> [--token <api_token>]
  python3 extract_office_doc.py --local <local_file_path>
"""

import argparse
import os
import sys
import tempfile
import urllib.request
import urllib.error
from pathlib import Path


def ensure_deps():
    """Install required libraries if missing."""
    import subprocess
    packages = {
        "pypdf": "pypdf",
        "docx": "python-docx",
        "pptx": "python-pptx",
        "openpyxl": "openpyxl",
    }
    for module, pkg in packages.items():
        try:
            __import__(module)
        except ImportError:
            print(f"[Installing {pkg}...]", file=sys.stderr)
            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", pkg,
                 "--break-system-packages", "-q"]
            )

ensure_deps()

from pypdf import PdfReader
import docx as python_docx
from pptx import Presentation
import openpyxl


def download_file(url: str, token: str = None) -> tuple:
    """Download file from URL. Returns (content_bytes, filename)."""
    req = urllib.request.Request(url)
    if token:
        req.add_header("Authorization", f"token {token}")

    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            content = resp.read()
            filename = url.split("?")[0].split("/")[-1]
            return content, filename
    except urllib.error.HTTPError as e:
        if e.code == 401:
            print(f"ERROR: Authentication failed (401). Check your API token.", file=sys.stderr)
        elif e.code == 403:
            print(f"ERROR: Access forbidden (403). You may not have permission.", file=sys.stderr)
        elif e.code == 404:
            print(f"ERROR: File not found (404). Check the URL and branch name.", file=sys.stderr)
        else:
            print(f"ERROR: HTTP {e.code} - {e.reason}", file=sys.stderr)
        sys.exit(1)
    except urllib.error.URLError as e:
        print(f"ERROR: Network error - {e.reason}", file=sys.stderr)
        sys.exit(1)


def extract_pdf(filepath: str) -> str:
    try:
        reader = PdfReader(filepath)
        total_pages = len(reader.pages)
        max_pages = min(total_pages, 20)

        lines = []
        lines.append(f"[PDF: {total_pages} pages total" +
                      (f", showing first {max_pages}" if total_pages > max_pages else "") + "]")

        for i, page in enumerate(reader.pages[:max_pages]):
            text = page.extract_text()
            if text and text.strip():
                lines.append(f"\n--- Page {i+1} ---")
                lines.append(text.strip())

        result = "\n".join(lines)
        text_content = "\n".join(lines[1:])
        if total_pages > 0 and len(text_content.strip()) < 50:
            return ("[PDF WARNING] Very little text extracted. This PDF may be a scanned image "
                    "without a text layer.\n" + result)
        return result
    except Exception as e:
        return f"[PDF ERROR] Failed to extract: {e}"


def extract_docx(filepath: str) -> str:
    try:
        doc = python_docx.Document(filepath)
        lines = []

        core_props = doc.core_properties
        if core_props.title:
            lines.append(f"[Document Title: {core_props.title}]")
        if core_props.author:
            lines.append(f"[Author: {core_props.author}]")
        lines.append("")

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "").strip()
                lines.append(f"\n{'#' * int(level) if level.isdigit() else '#'} {text}")
            else:
                lines.append(text)

        if doc.tables:
            lines.append("\n--- Tables ---")
            for t_idx, table in enumerate(doc.tables):
                lines.append(f"\n[Table {t_idx+1}]")
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(row_data))

        return "\n".join(lines)
    except Exception as e:
        return f"[DOCX ERROR] Failed to extract: {e}"


def extract_pptx(filepath: str) -> str:
    try:
        prs = Presentation(filepath)
        total_slides = len(prs.slides)
        max_slides = min(total_slides, 30)

        lines = []
        lines.append(f"[PowerPoint: {total_slides} slides total" +
                      (f", showing first {max_slides}" if total_slides > max_slides else "") + "]")

        for i in range(max_slides):
            slide = prs.slides[i]
            slide_texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if slide_texts:
                lines.append(f"\n[Slide {i+1}]")
                lines.extend(slide_texts)

        return "\n".join(lines)
    except Exception as e:
        return f"[PPTX ERROR] Failed to extract: {e}"


def extract_xlsx(filepath: str) -> str:
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
        lines = []
        lines.append(f"[Excel: {len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)}]")

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"\n--- Sheet: {sheet_name} ---")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if all(cell is None or str(cell).strip() == "" for cell in row):
                    continue
                row_str = " | ".join(str(c) if c is not None else "" for c in row)
                lines.append(row_str)
                row_count += 1
                if row_count >= 100:
                    lines.append("... (truncated, showing first 100 non-empty rows)")
                    break

        return "\n".join(lines)
    except Exception as e:
        return f"[XLSX ERROR] Failed to extract: {e}"


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".doc": extract_docx,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,
}


def extract(filepath: str, filename: str = None) -> str:
    name = filename or filepath
    ext = Path(name).suffix.lower()

    if ext not in EXTRACTORS:
        return (f"[UNSUPPORTED] File type '{ext}' is not supported.\n"
                f"Supported: {', '.join(EXTRACTORS.keys())}")

    return EXTRACTORS[ext](filepath)


def main():
    parser = argparse.ArgumentParser(description="Extract text from office documents in Gitea")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--url", help="Gitea raw file download URL")
    group.add_argument("--local", help="Local file path")
    parser.add_argument("--token", help="Gitea API token (for --url)")
    parser.add_argument("--filename", help="Override filename for type detection")

    args = parser.parse_args()

    if args.url:
        print(f"[Downloading from Gitea...]", file=sys.stderr)
        content, filename = download_file(args.url, args.token)
        override_name = args.filename or filename
        ext = Path(override_name).suffix.lower()
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            print(f"[Extracting content from: {override_name}]", file=sys.stderr)
            result = extract(tmp_path, override_name)
        finally:
            os.unlink(tmp_path)
    else:
        filepath = args.local
        if not os.path.exists(filepath):
            print(f"ERROR: File not found: {filepath}", file=sys.stderr)
            sys.exit(1)
        override_name = args.filename or filepath
        print(f"[Extracting content from local file: {override_name}]", file=sys.stderr)
        result = extract(filepath, override_name)

    print(result)


if __name__ == "__main__":
    main()